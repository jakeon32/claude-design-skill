"""
pptx_utils.py — python-pptx 공통 헬퍼 + DesignSystem

사용법:
    from pptx_utils import DesignSystem, PptxBuilder, DEFAULT_DS

    ds = DesignSystem(...)           # 커스텀 디자인 시스템
    b  = PptxBuilder(ds)
    prs = b.new_prs()
    sl  = b.blank(prs)
    b.add_txt(sl, ...)
    prs.save('out.pptx')
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from dataclasses import dataclass, field
from typing import Optional
import os

RECT  = 1
RRECT = 5
OVAL  = 9


# ─────────────────────────────────────────────────────
# PPTX Fidelity Patterns — LINE_SPACING_PT_TABLE (C-2)
# ─────────────────────────────────────────────────────
# PPTX runtime은 한국어 폰트 ascender+descender에 자동 leading을 추가하여
# ratio 기반 line_spacing이 시각적으로 1.15~1.25배 늘어남.
# size별로 'PPTX에서 시각 line-height가 X로 보이는 Pt 값'을 사전 캘리브레이션.
# v3.1 ferrari deck 6 핵심 슬라이드(S03/S04/S07/S08/S13/S14)에서 100% 검증 완료.
LINE_SPACING_PT_TABLE = [
    # (size_pt 임계값 ≥, line_spacing_pt)
    (96, 91),
    (84, 80),
    (72, 70),
    (56, 58),
    (40, 44),
    (28, 34),
    (20, 26),
    (16, 24),
    (14, 20),
    (12, 17),
    (10, 14),
    (8,  12),
]


# ─────────────────────────────────────────────────────
# DesignSystem
# ─────────────────────────────────────────────────────

@dataclass
class DesignSystem:
    """모든 디자인 토큰을 담는 데이터 클래스. 변경하고 싶은 값만 override하세요."""

    # 슬라이드 크기
    slide_w: float = 10.0   # inches
    slide_h: float = 5.625  # inches
    margin:  float = 0.7    # side margin
    gap:     float = 0.18   # default gap between elements

    # 색상 팔레트
    colors: dict = field(default_factory=lambda: {
        'bg':   RGBColor(0xFF, 0xFD, 0xF5),
        'fg':   RGBColor(0x1E, 0x29, 0x3B),
        'mute': RGBColor(0xF1, 0xF5, 0xF9),
        'mfg':  RGBColor(0x64, 0x74, 0x8B),
        'v':    RGBColor(0x8B, 0x5C, 0xF6),   # violet / accent
        'p':    RGBColor(0xF4, 0x72, 0xB6),   # pink
        'y':    RGBColor(0xFB, 0xBF, 0x24),   # yellow
        'm':    RGBColor(0x34, 0xD3, 0x99),   # mint
        'bdr':  RGBColor(0xE2, 0xE8, 0xF0),   # border
        'w':    RGBColor(0xFF, 0xFF, 0xFF),
        'w60':  RGBColor(0x99, 0xAA, 0xBB),
        'w80':  RGBColor(0xCC, 0xD5, 0xDE),
        'y_dk': RGBColor(0x92, 0x69, 0x0A),
        'g_dk': RGBColor(0x0D, 0x7A, 0x55),
    })

    # 폰트
    font_heading:      str = 'Pretendard'
    font_body:         str = 'Pretendard'
    font_heading_bold: str = 'Pretendard Black'   # 900 weight family name
    font_english:      str = 'Plus Jakarta Sans'  # 영문 전용


# 기본 디자인 시스템 (AI 인플루언서 테마)
DEFAULT_DS = DesignSystem()


# ─────────────────────────────────────────────────────
# PptxBuilder
# ─────────────────────────────────────────────────────

class PptxBuilder:
    """DesignSystem을 받아 슬라이드 요소를 생성하는 빌더."""

    def __init__(self, ds: DesignSystem = None):
        self.ds = ds or DEFAULT_DS
        self.W  = Inches(self.ds.slide_w)
        self.H  = Inches(self.ds.slide_h)
        self.M  = Inches(self.ds.margin)
        self.GAP = Inches(self.ds.gap)

    @property
    def C(self):
        return self.ds.colors

    @property
    def FH(self):  return self.ds.font_heading
    @property
    def FB(self):  return self.ds.font_body
    @property
    def FHB(self): return self.ds.font_heading_bold
    @property
    def FE(self):  return self.ds.font_english

    # ── Presentation ─────────────────────────────────

    def new_prs(self) -> Presentation:
        prs = Presentation()
        prs.slide_width  = self.W
        prs.slide_height = self.H
        return prs

    def blank(self, prs: Presentation):
        return prs.slides.add_slide(prs.slide_layouts[6])

    # ── Background / Fill ────────────────────────────

    def set_bg(self, slide, color: RGBColor):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def set_alpha(self, shape, pct: float):
        solid = shape.fill._xPr.find('.//' + qn('a:solidFill'))
        if solid is None or len(solid) == 0:
            return
        clr = solid[0]
        for a in clr.findall(qn('a:alpha')):
            clr.remove(a)
        el = etree.SubElement(clr, qn('a:alpha'))
        el.set('val', str(int(pct * 1000)))

    # ── Shapes ───────────────────────────────────────

    def add_rect(self, slide, x, y, w, h, fill: RGBColor,
                 line: Optional[RGBColor] = None, lw=None,
                 rounded=False, alpha=None, radius=None):
        lw = lw or Pt(1.5)
        s = slide.shapes.add_shape(RRECT if rounded else RECT, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if alpha is not None:
            self.set_alpha(s, alpha)
        if line:
            s.line.color.rgb = line
            s.line.width = lw
        else:
            s.line.fill.background()
        if rounded and radius is not None:
            self.set_radius(s, radius)
        self.no_shadow(s)
        return s

    def add_oval(self, slide, x, y, w, h, fill: RGBColor, alpha=100):
        s = slide.shapes.add_shape(OVAL, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.fill.background()
        if alpha < 100:
            self.set_alpha(s, alpha)
        self.no_shadow(s)
        return s

    def no_shadow(self, shape):
        spPr = shape._element.spPr
        for el in list(spPr):
            if el.tag in (qn('a:effectLst'), qn('a:effectDag')):
                spPr.remove(el)
        etree.SubElement(spPr, qn('a:effectLst'))

    def set_radius(self, shape, pct=8):
        prstGeom = shape._element.find('.//' + qn('a:prstGeom'))
        if prstGeom is None:
            return
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {int(pct * 1000)}')

    # ── Text ─────────────────────────────────────────

    def add_txt(self, slide, x, y, w, h, text, size=13, bold=False,
                color: Optional[RGBColor] = None,
                align=PP_ALIGN.LEFT, font=None, wrap=True, anchor=None):
        color = color or self.C['fg']
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        if anchor:
            tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font or self.FH
        return tb

    def add_mixed(self, slide, x, y, w, h, lines_parts,
                  size=36, wrap=True, align=PP_ALIGN.LEFT):
        """
        Multi-line heading with per-run colors.
        lines_parts: [ [(text, color), ...], ... ]
        """
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        first = True
        for line in lines_parts:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            for text, color in line:
                r = p.add_run()
                r.text = text
                r.font.size = Pt(size)
                r.font.bold = False
                r.font.color.rgb = color
                r.font.name = self.FHB
        return tb

    def add_lines(self, slide, x, y, w, h, lines, wrap=True):
        """
        Multi-paragraph text box.
        lines: [{'t': text, 's': size, 'b': bold, 'c': color, 'f': font,
                 'align': PP_ALIGN, 'sp': space_before_pt}, ...]
        """
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        first = True
        for ln in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = ln.get('align', PP_ALIGN.LEFT)
            if ln.get('sp'):
                p.space_before = Pt(ln['sp'])
            r = p.add_run()
            r.text = ln.get('t', '')
            r.font.size = Pt(ln.get('s', 13))
            r.font.bold = ln.get('b', False)
            r.font.color.rgb = ln.get('c', self.C['fg'])
            r.font.name = ln.get('f', self.FH)
        return tb

    # ── Composite widgets ────────────────────────────

    def pn(self, slide, n: int, total: int, dark_bg=False):
        """Page number (bottom-right)."""
        color = RGBColor(0x44, 0x55, 0x66) if dark_bg else RGBColor(0xC8, 0xD2, 0xDC)
        self.add_txt(slide,
                     self.W - Inches(1.3), self.H - Inches(0.36),
                     Inches(1.1), Inches(0.28),
                     f'{n:02d}  -  {total:02d}',
                     size=9, color=color, align=PP_ALIGN.RIGHT)

    def badge(self, slide, x, y, text, bg_color: RGBColor,
              text_color: Optional[RGBColor] = None):
        """Pill badge. Returns next y position."""
        tc = text_color or self.C['w']
        bw = max(Inches(1.5), Inches(len(text) * 0.085 + 0.4))
        bh = Inches(0.27)
        s = self.add_rect(slide, x, y, bw, bh, fill=bg_color, rounded=True)
        self.set_radius(s, 50)
        self.add_txt(slide, x, y, bw, bh, text.upper(), size=9, bold=True,
                     color=tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return y + bh + Inches(0.08)

    def card(self, slide, x, y, w, h, line=True):
        """White rounded card."""
        lc = self.C['bdr'] if line else None
        s = self.add_rect(slide, x, y, w, h, fill=self.C['w'], line=lc,
                          lw=Pt(1.5), rounded=True)
        self.set_radius(s, 6)
        return s

    def mixed_body(self, slide, x, y, w, h, pre, emph, post='',
                   emph_color: Optional[RGBColor] = None, size=11):
        """Body text with inline color emphasis (single paragraph)."""
        ec = emph_color or self.C['m']
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        for txt, clr, bold in [
            (pre,  self.C['w80'], False),
            (emph, ec,            True),
            (post, self.C['w80'], False),
        ]:
            if not txt:
                continue
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = clr
            r.font.name = self.FB
            r.font.bold = bold
        return tb

    # ── PPTX Fidelity Helpers (v3.1 검증) ────────────
    # ferrari deck v1→v3.1 에서 도출된 11개 패턴 중 코드 헬퍼화 가능한 4종 + 보조함수.
    # 디테일은 references/pptx-fidelity-patterns.md 참조.

    LINE_SPACING_PT_TABLE = LINE_SPACING_PT_TABLE

    def line_spacing_pt(self, size_pt, hero=False):
        """[C-2] font_size(pt) → line_spacing absolute Pt 값.

        PPTX runtime이 한국어 폰트 자동 leading을 추가하므로 ratio 대신
        Pt 절대값으로 line_spacing 지정. size별 사전 캘리브레이션.

        Parameters
        ----------
        size_pt : float
            font size in points
        hero : bool
            True면 한 단계 더 압축 (3% 추가) — 한 줄 hero 전용.
            두 줄 이상 hero는 F-1 패턴 사용 (multiline 인자).
        """
        for thr, sp in self.LINE_SPACING_PT_TABLE:
            if size_pt >= thr:
                base = sp
                break
        else:
            base = int(size_pt * 1.4)
        if hero and size_pt >= 56:
            # [v3.1 F-1] hero 명시 시 한 단계 더 압축 (3%)
            base = int(round(base * 0.97))
        return Pt(base)

    def _zero_margins(self, tf):
        """[D-3] textframe 내부 padding 모두 0 — 누적 오차 차단."""
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

    def _no_effect(self, shape):
        """그림자·effect 제거 — 헬퍼들이 공통으로 호출."""
        spPr = shape._element.spPr
        for el in list(spPr):
            if el.tag == qn('a:effectLst'):
                spPr.remove(el)
        etree.SubElement(spPr, qn('a:effectLst'))

    def add_hero_text(self, slide, x, y, w, h, lines_parts, size=80,
                      line_height=1.0, font=None, align=PP_ALIGN.LEFT,
                      hero=True, multiline=False):
        """[C-2 + C-3 + D-3 + F-1] hero/heading 텍스트박스.

        - C-2: size별 Pt 절대값 line_spacing 자동 적용 (size ≥ 40pt)
        - C-3: vertical_anchor=TOP 명시
        - D-3: zero margins 강제
        - F-1: multiline=True 시 line_spacing을 size × line_height로 직접 계산
               (테이블 압축 우회 — 두 줄 hero 마침표/다음 줄 충돌 방지)

        Parameters
        ----------
        lines_parts : list[list[tuple]]
            각 라인은 run tuple의 list. 다음 형식 지원:
              2-tuple (text, color) — bold=True 기본, size=base
              3-tuple (text, color, bold) — size=base
              4-tuple (text, color, bold, size_override) — 인라인 size 변경
        hero : bool
            True면 C-2 hero 압축 적용. 단일 줄 빅 데이터는 False 권장.
        multiline : bool
            True면 F-1 (size × line_height 직접 계산), 두 줄 이상 hero용.
        """
        font = font or self.FH
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        self._zero_margins(tf)                  # D-3
        tf.vertical_anchor = MSO_ANCHOR.TOP     # C-3

        # F-1 vs C-2 분기
        if multiline and size >= 40:
            ls = Pt(int(round(size * line_height)))
        elif size >= 40:
            ls = self.line_spacing_pt(size, hero=hero)
        else:
            ls = line_height  # ratio (소형 텍스트는 ratio도 OK)

        first = True
        for line in lines_parts:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.line_spacing = ls
            for part in line:
                # 2/3/4-tuple 지원
                if len(part) == 4:
                    text, color, bold, size_override = part
                elif len(part) == 3:
                    text, color, bold = part
                    size_override = None
                else:
                    text, color = part
                    bold = True
                    size_override = None
                r = p.add_run()
                r.text = text
                r.font.size = Pt(size_override if size_override is not None else size)
                r.font.bold = bold
                r.font.color.rgb = color
                r.font.name = font
        return tb

    def add_body_text(self, slide, x, y, w, h, text, size=12,
                      color: Optional[RGBColor] = None,
                      font=None, bold=False,
                      line_height=1.5, align=PP_ALIGN.LEFT, anchor=None):
        """[C-2 + D-3] 본문 textbox — size별 Pt 절대값 line_spacing 자동 적용.

        - size ≥ 40pt: hero=True 압축 (대형 단일 라인 호출자용)
        - 20 ≤ size < 40: hero=False (h2/h3 인라인)
        - size < 20: 본문 — Pt 절대값 그대로
        """
        color = color or self.C['fg']
        font = font or self.FB
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        self._zero_margins(tf)                  # D-3
        if anchor:
            tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        # C-2 size별 자동 분기
        if size >= 40:
            p.line_spacing = self.line_spacing_pt(size, hero=True)
        elif size >= 20:
            p.line_spacing = self.line_spacing_pt(size, hero=False)
        else:
            p.line_spacing = self.line_spacing_pt(size)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        return tb

    def add_para_bottom(self, slide, x, y, w, h, text, size=12,
                        color: Optional[RGBColor] = None, font=None,
                        bold=False, line_height=1.5, align=PP_ALIGN.LEFT):
        """[D-2 + D-3 + C-2] flex margin-top:auto → BOTTOM anchor textbox.

        parent box를 의미있게 키워야(예: h=28→44) anchor=BOTTOM 차이가 보임.
        """
        color = color or self.C['fg']
        font = font or self.FB
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        self._zero_margins(tf)                          # D-3
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM          # D-2 핵심
        p = tf.paragraphs[0]
        p.alignment = align
        # body 영역(작은 size) → Pt 절대값
        if size <= 18:
            p.line_spacing = self.line_spacing_pt(size)
        else:
            p.line_spacing = line_height
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        return tb

    def add_rule_line(self, slide, x, y, w, color: RGBColor, weight_pt=0.75):
        """[E-1] 1px rule → 얇은 직사각형 (rect h=0.75pt).

        ※ MSO_CONNECTOR_TYPE 사용 절대 금지. PowerPoint Application.Open()이
           cxnSp + schemeClr 조합을 거부하는 사례가 있음. rect로 안정화.
        """
        h_emu = Pt(weight_pt)
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h_emu)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        self._no_effect(s)
        return s

    def add_image_cover(self, slide, path, x, y, w, h):
        """[B-1] object-fit:cover 시뮬레이션.

        컨테이너 비율과 이미지 비율을 비교 → crop_left/right/top/bottom 적용.
        path가 없거나 PIL 미설치 시 단색 fallback rect.
        """
        if not os.path.exists(path):
            s = slide.shapes.add_shape(RECT, x, y, w, h)
            s.fill.solid()
            s.fill.fore_color.rgb = self.C.get('fg', RGBColor(0x18, 0x18, 0x18))
            s.line.fill.background()
            return None

        try:
            from PIL import Image
        except ImportError:
            return slide.shapes.add_picture(path, x, y, width=w, height=h)

        try:
            with Image.open(path) as im:
                iw, ih = im.size
        except Exception:
            return slide.shapes.add_picture(path, x, y, width=w, height=h)

        img_ratio = iw / ih
        cw_emu = int(w)
        ch_emu = int(h)
        cont_ratio = cw_emu / ch_emu

        pic = slide.shapes.add_picture(path, x, y, width=w, height=h)

        if img_ratio > cont_ratio:
            # 이미지가 더 가로형 → 좌우 crop
            keep = cont_ratio / img_ratio
            crop = (1 - keep) / 2
            try:
                pic.crop_left = crop
                pic.crop_right = crop
            except AttributeError:
                self._set_picture_crop(pic, l=crop, r=crop)
        elif img_ratio < cont_ratio:
            # 이미지가 더 세로형 → 위아래 crop
            keep = img_ratio / cont_ratio
            crop = (1 - keep) / 2
            try:
                pic.crop_top = crop
                pic.crop_bottom = crop
            except AttributeError:
                self._set_picture_crop(pic, t=crop, b=crop)
        return pic

    def _set_picture_crop(self, pic, l=0, r=0, t=0, b=0):
        """python-pptx <0.6.21 fallback — XML srcRect 직접 조작."""
        blipFill = pic._element.blipFill
        srcRect = blipFill.find(qn('a:srcRect'))
        if srcRect is None:
            srcRect = etree.SubElement(blipFill, qn('a:srcRect'))
            blip = blipFill.find(qn('a:blip'))
            if blip is not None:
                blip.addnext(srcRect)
        if l: srcRect.set('l', str(int(l * 100000)))
        if r: srcRect.set('r', str(int(r * 100000)))
        if t: srcRect.set('t', str(int(t * 100000)))
        if b: srcRect.set('b', str(int(b * 100000)))


# ─────────────────────────────────────────────────────
# build_pptx — 범용 진입점
# ─────────────────────────────────────────────────────

def build_pptx(slide_builders, output_path: str,
               ds: DesignSystem = None) -> Presentation:
    """
    Parameters
    ----------
    slide_builders : list[callable]
        각 함수는 (prs, builder) → slide 시그니처.
        예: [cover_slide, intro_slide, ...]
    output_path : str
        저장 경로 (.pptx)
    ds : DesignSystem, optional
        디자인 시스템. None이면 DEFAULT_DS 사용.

    Returns
    -------
    Presentation
    """
    b = PptxBuilder(ds)
    prs = b.new_prs()
    for fn in slide_builders:
        fn(prs, b)
    prs.save(output_path)
    print(f'Saved: {output_path}  ({len(prs.slides)} slides)')
    return prs
